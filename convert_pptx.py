"""
PowerPointファイル（PPTX）を自動変換するスクリプト

入力PPT: 表紙 → 問題ページ群 → 解答1枚
出力PPT: 表紙 → 問題ページ群 → 解答を埋め込んだ問題ページ群 → 解答1枚
"""

import os
import re
import sys
import shutil
import io
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn


def extract_answers(answer_slide):
    """
    解答ページから大問番号と解答文字列の辞書を生成
    
    Args:
        answer_slide: 解答ページのスライドオブジェクト
        
    Returns:
        dict: {大問番号: 解答文字列} の辞書
        例: {1: "38.608㎠", 2: "52.685㎠", 3: "36", ...}
    """
    # スライド内のすべてのテキストを連結
    all_text = ""
    for shape in answer_slide.shapes:
        if hasattr(shape, "text"):
            all_text += shape.text + "\n"
    
    # 正規表現で解答を抽出
    # パターン: 大問\s*([0-9]+)\s*([^\n\r]+)
    pattern = r"大問\s*([0-9]+)\s*([^\n\r]+)"
    matches = re.findall(pattern, all_text)
    
    answers = {}
    for match in matches:
        question_num = int(match[0])
        answer_text = match[1].strip()
        answers[question_num] = answer_text
    
    return answers


def extract_text_from_shape(shape):
    """
    シェイプからテキストを取得（GROUP内の子シェイプも再帰的に探索）
    
    Args:
        shape: シェイプオブジェクト
        
    Returns:
        str: テキスト（見つからない場合は空文字列）
    """
    text = ""
    try:
        # 直接テキストを取得
        if hasattr(shape, "text") and shape.text:
            text = shape.text.strip()
        elif hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
        
        # GROUPシェイプの場合は子シェイプも探索
        if not text and hasattr(shape, "shapes"):
            for child_shape in shape.shapes:
                child_text = extract_text_from_shape(child_shape)
                if child_text:
                    text = child_text
                    break
    except:
        pass
    
    return text


def get_font_size_from_shape(shape):
    """
    シェイプからフォントサイズを取得（GROUP内の子シェイプも再帰的に探索）
    
    Args:
        shape: シェイプオブジェクト
        
    Returns:
        float: フォントサイズ（ポイント単位）、見つからない場合はNone
    """
    font_size = None
    try:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            if len(shape.text_frame.paragraphs) > 0:
                paragraph = shape.text_frame.paragraphs[0]
                if paragraph.runs:
                    # 最初のrunのフォントサイズを取得
                    font = paragraph.runs[0].font
                    if font and font.size:
                        font_size = font.size.pt  # ポイント単位で取得
        
        # GROUPシェイプの場合は子シェイプも探索
        if font_size is None and hasattr(shape, "shapes"):
            for child_shape in shape.shapes:
                child_font_size = get_font_size_from_shape(child_shape)
                if child_font_size is not None:
                    font_size = child_font_size
                    break
    except:
        pass
    
    return font_size


def extract_question_number_candidates(slide, prs, debug=False):
    """
    問題ページから大問番号の候補を抽出
    スライドを左右に分割し、それぞれの領域の左上から大問番号を探す
    GROUPシェイプ内の子シェイプも再帰的に探索
    
    Args:
        slide: 問題ページのスライドオブジェクト
        prs: Presentationオブジェクト（スライドサイズ取得用）
        debug: デバッグ情報を出力するかどうか
        
    Returns:
        list: 大問番号の候補リスト（(優先度, 距離, 大問番号, テキスト, 領域)のタプルのリスト）
        領域は 'left' または 'right'
    """
    # スライドのサイズを取得
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # スライドを左右に分割
    left_half_end = slide_width / 2
    right_half_start = slide_width / 2
    
    # 各領域の左上の範囲を定義（各領域の左上15%以内に絞る）
    area_width = slide_width / 2
    left_threshold = area_width * 0.15
    top_threshold = slide_height * 0.15
    
    candidates = []
    debug_shapes = []  # デバッグ用
    
    def process_shape(shape, parent_left=0, parent_top=0):
        """シェイプを再帰的に処理"""
        try:
            # シェイプの位置を取得（GROUP内の場合は親からの相対位置、直接の場合は絶対位置）
            if hasattr(shape, 'left') and hasattr(shape, 'top'):
                # python-pptxでは、GROUP内の子シェイプの位置は親からの相対位置として取得される
                # 親の位置を加算する必要がある
                shape_left = parent_left + shape.left
                shape_top = parent_top + shape.top
            else:
                shape_left = parent_left
                shape_top = parent_top
            
            # 左半分の領域か右半分の領域かを判定
            area = None
            relative_left = None
            
            if shape_left < left_half_end:
                # 左半分
                area = 'left'
                relative_left = shape_left
            elif shape_left >= right_half_start:
                # 右半分
                area = 'right'
                relative_left = shape_left - right_half_start
            
            if area is None:
                # GROUP内の子シェイプを探索
                if hasattr(shape, "shapes"):
                    for child_shape in shape.shapes:
                        process_shape(child_shape, shape_left, shape_top)
                return
            
            # 各領域の左上15%以内か確認
            if relative_left <= left_threshold and shape_top <= top_threshold:
                # テキストを取得（GROUP内も再帰的に探索）
                text = extract_text_from_shape(shape)
                
                # デバッグ情報を記録（短いテキストのみ）
                if debug and len(text) <= 10:
                    shape_type = getattr(shape, 'shape_type', 'unknown')
                    debug_shapes.append({
                        'area': area,
                        'text': text,
                        'left': shape_left,
                        'top': shape_top,
                        'relative_left': relative_left,
                        'shape_type': shape_type
                    })
                
                # 大問番号は通常、短いテキスト（1-3文字程度）で数字のみ
                # 長いテキスト（問題文など）は除外
                if len(text) > 5:
                    # GROUP内の子シェイプを探索
                    if hasattr(shape, "shapes"):
                        for child_shape in shape.shapes:
                            process_shape(child_shape, shape_left, shape_top)
                    return
                
                # テキストが数字のみの場合のみを対象
                if not text.isdigit():
                    # GROUP内の子シェイプを探索
                    if hasattr(shape, "shapes"):
                        for child_shape in shape.shapes:
                            process_shape(child_shape, shape_left, shape_top)
                    return
                
                # 大問番号を取得
                question_num = int(text)
                
                # 大問番号は通常1-20程度なので、それ以外は除外
                if question_num < 1 or question_num > 20:
                    return
                
                # 位置でソートするため、距離を計算（各領域の左上からの距離）
                distance = (relative_left ** 2 + shape_top ** 2) ** 0.5
                
                # テキストボックスかどうかを判定
                is_textbox = False
                try:
                    if hasattr(shape, 'shape_type'):
                        if shape.shape_type in [1, 17]:  # MSO_SHAPE_TYPE.AUTO_SHAPE または TEXT_BOX
                            is_textbox = True
                except:
                    pass
                
                # 優先度: テキストボックス > 数字のみのテキスト > その他
                priority = 2 if is_textbox else (1 if text.isdigit() else 0)
                
                candidates.append((priority, distance, question_num, text, area))
            else:
                # GROUP内の子シェイプを探索
                if hasattr(shape, "shapes"):
                    for child_shape in shape.shapes:
                        process_shape(child_shape, shape_left, shape_top)
        except Exception as e:
            if debug:
                print(f"    シェイプ処理エラー: {e}")
    
    # すべてのシェイプを処理
    for shape in slide.shapes:
        process_shape(shape)
    
    # デバッグ情報を出力
    if debug and debug_shapes:
        print(f"    デバッグ: 左上領域内の短いテキストシェイプ {len(debug_shapes)}個")
        for ds in debug_shapes[:10]:  # 最初の10個を表示
            print(f"      {ds['area']}: '{ds['text']}' (長さ: {len(ds['text'])}, 位置: {ds['left']:.0f}, {ds['top']:.0f}, タイプ: {ds['shape_type']})")
    
    # 優先度と距離でソート（優先度が高い順、距離が小さい順）
    candidates.sort(key=lambda x: (-x[0], x[1]))
    return candidates


def match_question_numbers(question_slides, prs, answer_numbers):
    """
    各問題ページから抽出した大問番号候補と、解答ページの大問番号を照合
    各スライドには左右2つの問題があり、それぞれに大問番号がある
    
    Args:
        question_slides: 問題ページのスライドリスト
        prs: Presentationオブジェクト
        answer_numbers: 解答ページから抽出した大問番号のセット
        
    Returns:
        list: 各スライドの(左側の大問番号, 右側の大問番号)のタプルのリスト
    """
    # 各問題ページから大問番号の候補を取得（左右両方）
    all_candidates = []
    for idx, slide in enumerate(question_slides):
        candidates = extract_question_number_candidates(slide, prs, debug=True)
        all_candidates.append(candidates)
        if not candidates:
            print(f"  スライド{idx + 1}: 候補が見つかりませんでした")
    
    # 各スライドから左右の大問番号を抽出
    slide_question_numbers = []
    all_question_numbers = []  # 検証用
    
    for idx, candidates in enumerate(all_candidates):
        # 左右の領域からそれぞれ最適な候補を選択
        left_candidates = [c for c in candidates if c[4] == 'left']  # 領域が'left'のもの
        right_candidates = [c for c in candidates if c[4] == 'right']  # 領域が'right'のもの
        
        # 左側の大問番号を選択
        left_num = None
        if left_candidates:
            for priority, distance, question_num, text, area in left_candidates:
                if question_num in answer_numbers:
                    left_num = question_num
                    print(f"  スライド{idx + 1}（左側）: 大問番号 {question_num} を抽出 (テキスト: '{text}')")
                    break
            if left_num is None:
                left_num = left_candidates[0][2]
                print(f"警告: スライド{idx + 1}（左側）の大問番号候補 {left_num} は解答ページに存在しません。")
        else:
            print(f"警告: スライド{idx + 1}（左側）から大問番号の候補が見つかりませんでした。")
        
        # 右側の大問番号を選択
        right_num = None
        if right_candidates:
            for priority, distance, question_num, text, area in right_candidates:
                if question_num in answer_numbers:
                    right_num = question_num
                    print(f"  スライド{idx + 1}（右側）: 大問番号 {question_num} を抽出 (テキスト: '{text}')")
                    break
            if right_num is None:
                right_num = right_candidates[0][2]
                print(f"警告: スライド{idx + 1}（右側）の大問番号候補 {right_num} は解答ページに存在しません。")
        else:
            print(f"警告: スライド{idx + 1}（右側）から大問番号の候補が見つかりませんでした。")
        
        # 左右の大問番号のタプルを追加
        slide_question_numbers.append((left_num, right_num))
        if left_num is not None:
            all_question_numbers.append(left_num)
        if right_num is not None:
            all_question_numbers.append(right_num)
    
    # 検証: 1から順番に並んでいるか、合計が一致するか確認
    sorted_numbers = sorted(all_question_numbers)
    expected_sum = sum(answer_numbers)
    actual_sum = sum(all_question_numbers)
    
    print(f"\n大問番号の検証:")
    print(f"  抽出された大問番号: {all_question_numbers}")
    print(f"  解答ページの大問番号: {sorted(answer_numbers)}")
    print(f"  合計: 抽出={actual_sum}, 解答={expected_sum}")
    
    if sorted_numbers == list(range(1, len(sorted_numbers) + 1)) and len(sorted_numbers) == len(answer_numbers):
        print(f"  ✓ 1から順番に並んでいます")
    else:
        print(f"  警告: 1から順番に並んでいません")
    
    if actual_sum == expected_sum:
        print(f"  ✓ 合計が一致しています")
    else:
        print(f"  警告: 合計が一致しません")
    
    return slide_question_numbers


def duplicate_slide_complete(prs, source_slide):
    """
    スライドを完全複製する（画像やコメントも保持）
    Ctrl+C/Ctrl+Vのようにスライド全体を完全にコピー
    GROUPシェイプ内のすべての要素（テキスト、線、図形、画像など）も含めてコピー
    
    Args:
        prs: Presentationオブジェクト
        source_slide: 複製元のスライド
        
    Returns:
        複製されたスライドオブジェクト
    """
    import io
    
    # 新しいスライドを作成（同じレイアウトを使用）
    dest_slide = prs.slides.add_slide(source_slide.slide_layout)
    
    # デフォルトのプレースホルダーを削除（新規作成時に自動で入るテキストボックスを削除）
    for shape in list(dest_slide.shapes):
        try:
            # プレースホルダーかどうかを確認
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                sp = shape.element
                sp.getparent().remove(sp)
            else:
                # プレースホルダーでない場合も削除（レイアウトから来る要素）
                sp = shape.element
                sp.getparent().remove(sp)
        except:
            # 削除できない場合はスキップ
            continue
    
    # 元のスライドのすべてのシェイプをコピー
    for shape in source_slide.shapes:
        try:
            # 画像の場合は特別な処理
            if hasattr(shape, 'image') and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                # 画像データを取得
                image_stream = io.BytesIO(shape.image.blob)
                # 新しいスライドに画像を追加
                dest_slide.shapes.add_picture(
                    image_stream,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )
            else:
                # その他のシェイプ（GROUPシェイプを含む）はXMLをコピー
                # GROUPシェイプのXMLには子シェイプも含まれているため、これでコピーされる
                shape_xml = shape.element
                new_shape_xml = deepcopy(shape_xml)
                
                # GROUPシェイプの場合、XML内の画像リレーションシップのrIdを更新
                if hasattr(shape, 'shape_type') and shape.shape_type == 6:  # GROUP
                    # XML内のすべての画像リレーションシップを探してコピー
                    try:
                        # XML要素内のすべてのrId参照を探す（名前空間を考慮）
                        def find_and_update_rids(element, dest_slide, source_slide):
                            """XML要素内のrId参照を探して更新"""
                            rId_map = {}  # 元のrId -> 新しいrIdのマッピング
                            
                            # ソーススライドのリレーションシップを取得
                            source_rels = source_slide.part.rels
                            
                            # まず、XML内で使用されているすべてのrIdを収集
                            used_rIds = set()
                            ns_embed = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                            
                            for elem in element.iter():
                                # r:embed属性を確認
                                if ns_embed in elem.attrib:
                                    used_rIds.add(elem.attrib[ns_embed])
                                # その他のrId属性も確認
                                for attr_name, attr_value in list(elem.attrib.items()):
                                    if isinstance(attr_value, str) and attr_value.startswith('rId'):
                                        used_rIds.add(attr_value)
                            
                            # デバッグ: 見つかったrIdを表示
                            if used_rIds:
                                print(f"  デバッグ: GROUPシェイプ内で見つかったrId: {used_rIds}")
                            
                            # 使用されているrIdの画像リレーションシップをコピー
                            for rId in used_rIds:
                                # スライドレベルのリレーションシップを確認
                                if rId in source_rels:
                                    try:
                                        rel = source_rels[rId]
                                        rel_type = rel.reltype
                                        
                                        # 画像のリレーションシップの場合（PNG/JPEG/GIFなど）
                                        if "image" in rel_type or "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" in rel_type:
                                            # 画像データ（blob）を取得
                                            try:
                                                from pptx.parts.image import ImagePart
                                                image_blob = None
                                                
                                                # target_partの型を確認
                                                target_part = rel.target_part
                                                
                                                # ImagePartの場合、blobを取得
                                                if isinstance(target_part, ImagePart):
                                                    image_blob = target_part.blob
                                                else:
                                                    # target_partがImagePartでない場合、_targetから画像パートを取得
                                                    try:
                                                        target_path = rel._target
                                                        # packageから画像パートを取得
                                                        # target_pathが相対パスの場合、package.get_partで取得できる
                                                        image_part = source_slide.part.package.get_part(target_path)
                                                        if hasattr(image_part, 'blob'):
                                                            image_blob = image_part.blob
                                                        elif hasattr(image_part, 'image'):
                                                            # imageプロパティがある場合（一部のパートタイプ）
                                                            image_blob = image_part.image.blob
                                                        else:
                                                            # パートから直接読み込む
                                                            image_blob = image_part._blob if hasattr(image_part, '_blob') else None
                                                    except Exception as e:
                                                        # _targetから取得できない場合、target_partから直接取得を試みる
                                                        if hasattr(target_part, 'blob'):
                                                            image_blob = target_part.blob
                                                        elif hasattr(target_part, 'image'):
                                                            image_blob = target_part.image.blob
                                                        else:
                                                            print(f"  デバッグ: 画像blob取得エラー (rId: {rId}): {e}")
                                                
                                                if image_blob is None:
                                                    print(f"  デバッグ: 警告: 画像blobが取得できませんでした (rId: {rId})")
                                                    continue
                                                
                                                print(f"  デバッグ: 画像リレーションシップ {rId} を処理中 (タイプ: {rel_type}, サイズ: {len(image_blob)} bytes)")
                                                
                                                # 既に同じ画像データのリレーションシップが存在するか確認
                                                # スライドレベルのリレーションシップを確認
                                                already_exists = False
                                                existing_rId = None
                                                for dest_rel in dest_slide.part.rels:
                                                    try:
                                                        if dest_slide.part.rels[dest_rel].reltype == rel_type:
                                                            dest_target_part = dest_slide.part.rels[dest_rel].target_part
                                                            if isinstance(dest_target_part, ImagePart):
                                                                dest_blob = dest_target_part.blob
                                                                if dest_blob == image_blob:
                                                                    already_exists = True
                                                                    existing_rId = dest_rel
                                                                    print(f"  デバッグ: 既存のリレーションシップ {dest_rel} を使用")
                                                                    break
                                                    except:
                                                        continue
                                                
                                                if not already_exists:
                                                    # 画像パートは既にパッケージ内に存在している（shutil.copy2でコピー済み）
                                                    # 元のスライドのリレーションシップから画像パートのパスを取得して、パッケージから直接取得
                                                    try:
                                                        # 元のスライドの画像パートのパスを取得
                                                        source_image_path = rel._target
                                                        
                                                        # パッケージから同じパスで画像パートを取得
                                                        # （パッケージはコピーされているので、同じパスで画像パートが存在するはず）
                                                        found_image_part = None
                                                        try:
                                                            found_image_part = dest_slide.part.package.get_part(source_image_path)
                                                            print(f"  デバッグ: パッケージ内で既存の画像パートを発見 (パス: {source_image_path})")
                                                        except Exception as path_error:
                                                            # パスで取得できない場合、get_or_add_image_partを使用
                                                            print(f"  デバッグ: パスで取得できなかったため、get_or_add_image_partを使用: {path_error}")
                                                            image_stream = io.BytesIO(image_blob)
                                                            image_stream.seek(0)
                                                            found_image_part = dest_slide.part.package.get_or_add_image_part(image_stream)
                                                        
                                                        # 新しいリレーションシップを作成（スライドレベル）
                                                        # relate_toメソッドを使用
                                                        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
                                                        new_rId = dest_slide.part.relate_to(found_image_part, RT.IMAGE)
                                                        rId_map[rId] = new_rId
                                                        print(f"  デバッグ: 新しいリレーションシップ {new_rId} を作成 (元のrId: {rId})")
                                                    except Exception as e:
                                                        print(f"  デバッグ: 画像パート追加エラー (rId: {rId}): {e}")
                                                        print(f"  デバッグ: エラーの型: {type(e)}")
                                                        import traceback
                                                        print(f"  デバッグ: トレースバック:\n{traceback.format_exc()}")
                                                        # エラーが発生しても処理を続行
                                                        continue
                                                else:
                                                    rId_map[rId] = existing_rId
                                                    print(f"  デバッグ: rIdマッピング {rId} -> {existing_rId}")
                                            except Exception as e:
                                                # blob取得に失敗した場合はスキップ
                                                print(f"  デバッグ: 画像blob取得エラー (rId: {rId}): {e}")
                                                pass
                                    except Exception as e:
                                        print(f"  デバッグ: リレーションシップ処理エラー (rId: {rId}): {e}")
                                        pass
                            
                            # XML内のrId参照を更新
                            updated_count = 0
                            for elem in element.iter():
                                # r:embed属性を更新
                                if ns_embed in elem.attrib:
                                    old_rId = elem.attrib[ns_embed]
                                    if old_rId in rId_map:
                                        elem.attrib[ns_embed] = rId_map[old_rId]
                                        updated_count += 1
                                        print(f"  デバッグ: r:embed属性を更新 {old_rId} -> {rId_map[old_rId]}")
                                
                                # その他のrId属性も更新
                                for attr_name, attr_value in list(elem.attrib.items()):
                                    if isinstance(attr_value, str) and attr_value.startswith('rId'):
                                        if attr_value in rId_map:
                                            elem.attrib[attr_name] = rId_map[attr_value]
                                            updated_count += 1
                                            print(f"  デバッグ: {attr_name}属性を更新 {attr_value} -> {rId_map[attr_value]}")
                            
                            if updated_count == 0 and used_rIds:
                                print(f"  デバッグ: 警告: rIdマッピングが作成されませんでした (使用されたrId: {used_rIds})")
                            elif updated_count > 0:
                                print(f"  デバッグ: {updated_count}個のrId参照を更新しました")
                        
                        find_and_update_rids(new_shape_xml, dest_slide, source_slide)
                    except Exception as e:
                        # XML更新に失敗した場合は元のXMLを使用
                        pass
                
                dest_slide.shapes._spTree.append(new_shape_xml)
        except Exception as e:
            # コピーできないシェイプはスキップ
            continue
    
    # リレーションシップをコピー（画像などのリンクを保持）
    # 注意: リレーションシップのコピーは複雑なため、必要に応じて後で処理
    # XMLをコピーした時点で、多くのリレーションシップは既に参照されている
    try:
        for rel in source_slide.part.rels:
            if "notesSlide" not in source_slide.part.rels[rel].reltype:
                try:
                    # 既存のリレーションシップを確認
                    rel_target = source_slide.part.rels[rel]._target
                    rel_type = source_slide.part.rels[rel].reltype
                    
                    # 既に同じターゲットのリレーションシップが存在するか確認
                    already_exists = False
                    for dest_rel in dest_slide.part.rels:
                        if dest_slide.part.rels[dest_rel]._target == rel_target:
                            already_exists = True
                            break
                    
                    if not already_exists:
                        # 新しいリレーションシップを追加
                        dest_slide.part.rels.add_relationship(
                            rel_type,
                            rel_target,
                            rel.rId
                        )
                except:
                    # リレーションシップの追加に失敗した場合はスキップ
                    # （既に存在する場合や、rIdの重複など）
                    pass
    except:
        # リレーションシップのコピーに失敗しても続行
        pass
    
    return dest_slide


def find_answer_textbox(slide, prs, position='right'):
    """
    解答欄のテキストボックスを特定する
    
    解答欄の特徴：
    - 各問題の右下にある
    - 横向きの直線のオブジェクトがテキストボックス直下にある
    - さらに線の下に「点」というテキストが必ずある
    
    Args:
        slide: 対象のスライドオブジェクト
        prs: Presentationオブジェクト（スライドサイズ取得用）
        position: 'left' または 'right'（左右どちらの領域を探すか）
        
    Returns:
        解答欄のテキストボックスシェイプ、見つからない場合はNone
    """
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # 左右どちらの領域かを判定
    # 左側: 0 ～ slide_width * 0.62 (中央より少し左まで、左側の「点」が中央下の少し左にあるため、幅を広げる)
    # 右側: slide_width * 0.38 ～ slide_width (中央より少し右から、右側の「点」と重複しないように)
    # 注意: 左側の「点」は右端が中央付近にあるため、左端は画面左端付近まで行く場合がある
    left_area_end = slide_width * 0.62
    right_area_start = slide_width * 0.38
    
    # 「点」というテキストを持つシェイプを探す
    ten_shapes = []  # 「点」のシェイプのリスト
    all_ten_shapes = []  # すべての「点」のシェイプ（デバッグ用）
    
    def get_group_offsets(group_shape):
        """GROUPシェイプのa:offとa:chOffを取得"""
        try:
            from pptx.oxml.ns import qn
            grpSpPr = group_shape._element.find(qn('p:grpSpPr'))
            if grpSpPr is None:
                return None, None, None, None
            
            xfrm = grpSpPr.find(qn('a:xfrm'))
            if xfrm is None:
                return None, None, None, None
            
            # a:off（グループの位置）
            off = xfrm.find(qn('a:off'))
            if off is not None:
                group_off_x = int(off.get('x', 0))
                group_off_y = int(off.get('y', 0))
            else:
                group_off_x = group_off_y = 0
            
            # a:chOff（子シェイプの座標系の原点）
            chOff = xfrm.find(qn('a:chOff'))
            if chOff is not None:
                group_chOff_x = int(chOff.get('x', 0))
                group_chOff_y = int(chOff.get('y', 0))
            else:
                group_chOff_x = group_chOff_y = 0
            
            return group_off_x, group_off_y, group_chOff_x, group_chOff_y
        except:
            return None, None, None, None
    
    def process_shape(shape, parent_left=0, parent_top=0, depth=0, group_off_x=None, group_chOff_x=None, group_off_y=None, group_chOff_y=None, parent_group_shape=None):
        """シェイプを再帰的に処理（GROUPシェイプ内の座標を絶対座標に変換）"""
        try:
            # シェイプの位置を取得
            is_group = hasattr(shape, 'shape_type') and shape.shape_type == 6  # 6 = MSO_SHAPE_TYPE.GROUP
            # 現在のGROUPシェイプを記録（GROUPシェイプ自体の場合は自分自身、子シェイプの場合は親のGROUPシェイプ）
            current_group_shape = shape if is_group else parent_group_shape
            shape_relative_left = 0
            shape_relative_top = 0
            
            if hasattr(shape, 'left') and hasattr(shape, 'top'):
                shape_relative_left = shape.left
                shape_relative_top = shape.top
                
                # GROUPシェイプ内の子シェイプの場合、a:offとa:chOffを考慮して絶対座標に変換
                if group_off_x is not None and group_chOff_x is not None:
                    # absolute_x = child_shape.left + group_off_x - group_chOff_x
                    # absolute_y = child_shape.top + group_off_y - group_chOff_y
                    shape_left = shape_relative_left + group_off_x - group_chOff_x
                    shape_top = shape_relative_top + group_off_y - group_chOff_y
                else:
                    # 通常のシェイプまたはGROUPシェイプ自体の場合
                    shape_left = parent_left + shape_relative_left
                    shape_top = parent_top + shape_relative_top
            else:
                shape_left = parent_left
                shape_top = parent_top
            
            # テキストを取得（すべてのシェイプで確認）
            text = extract_text_from_shape(shape)
            
            # デバッグ: 短いテキスト（1-3文字）を持つシェイプを記録（左側の領域のみ）
            if position == 'left' and len(text) <= 3 and text.strip():
                shape_type = shape.shape_type if hasattr(shape, 'shape_type') else 'unknown'
                print(f"  デバッグ: 左側領域の短いテキスト: '{text}' (位置: {shape_left:.0f}, {shape_top:.0f}, タイプ: {shape_type})")
            
            # 「点」のみのテキストを持つシェイプを探す（すべての「点」を記録）
            # 「点」のみ = テキストが「点」だけで、他に何も含まれていない（空白を除く）
            text_stripped = text.strip()
            # 「点」のみかどうかを厳密に判定（全角・半角の空白を除いて「点」1文字のみ）
            is_ten_only = (text_stripped == "点" and len(text_stripped) == 1)
            
            if is_ten_only:
                # 解答欄の「点」の特徴で判定
                # 1. フォントサイズが8ptの「点」のみを対象（解答欄のテキストボックスは9pt）
                # 2. スライドの半分（50%）よりも下の位置のテキストボックスのみを対象
                font_size = get_font_size_from_shape(shape)
                
                # フォントサイズが8ptで、かつスライドの下半分にある「点」のみを対象
                is_font_8pt = (font_size == 8.0)
                is_bottom_half = (shape_top > slide_height / 2) if slide_height > 0 else False
                is_answer_ten = is_font_8pt and is_bottom_half
                
                if is_answer_ten:
                    shape_width = shape.width if hasattr(shape, 'width') else 0
                    shape_height = shape.height if hasattr(shape, 'height') else 0
                    shape_right = shape_left + shape_width  # 右端座標
                    
                    # 左右の分類は後で中央基準で行うため、ここでは一時的にNone
                    area = None
                    
                    all_ten_shapes.append({
                        'shape': shape,
                        'left': shape_left,
                        'top': shape_top,
                        'right': shape_right,  # 右端座標を追加
                        'width': shape_width,
                        'height': shape_height,
                        'area': area,
                        'position_ratio': shape_left / slide_width if slide_width > 0 else 0,
                        'text': text,  # デバッグ用に元のテキストも保存
                        'text_stripped': text_stripped,  # デバッグ用
                        'font_size': font_size,  # デバッグ用
                        'parent_group_shape': current_group_shape  # 親のGROUPシェイプを記録
                    })
                    top_ratio = (shape_top / slide_height * 100) if slide_height > 0 else 0
                    right_ratio = (shape_right / slide_width * 100) if slide_width > 0 else 0
                    print(f"  デバッグ: 「点」のみを発見（解答欄、フォント8pt、下半分）")
                    print(f"    位置: 左端={shape_left:.0f}, 右端={shape_right:.0f}, 上端={shape_top:.0f}, 下端={shape_top + shape_height:.0f}")
                    print(f"    サイズ: 幅={shape_width:.0f}, 高さ={shape_height:.0f}")
                    print(f"    比率: 右端={right_ratio:.1f}%, 上端={top_ratio:.1f}%, フォントサイズ={font_size}pt")
                    print(f"    座標情報: 親座標=({parent_left:.0f}, {parent_top:.0f}), 相対座標=({shape_relative_left:.0f}, {shape_relative_top:.0f}), 深度={depth}, GROUP内={is_group}")
                else:
                    reason = []
                    if not is_font_8pt:
                        reason.append(f"フォントサイズ不一致({font_size}pt, 期待値: 8pt)")
                    if not is_bottom_half:
                        top_ratio = (shape_top / slide_height * 100) if slide_height > 0 else 0
                        reason.append(f"上半分({top_ratio:.1f}%)")
                    reason_str = ", ".join(reason) if reason else "不明"
                    print(f"  デバッグ: 「点」のみを除外 ({reason_str}) (位置: {shape_left:.0f}, {shape_top:.0f})")
                
                # 指定された領域の「点」のみをten_shapesに追加
                if area == position:
                    ten_shapes.append({
                        'shape': shape,
                        'left': shape_left,
                        'top': shape_top,
                        'width': shape.width if hasattr(shape, 'width') else 0,
                        'height': shape.height if hasattr(shape, 'height') else 0
                    })
                    print(f"  デバッグ: 対象領域の「点」のみを発見 ({position}側, 位置: {shape_left:.0f}, {shape_top:.0f})")
            
            # GROUPシェイプの場合は子シェイプも探索
            if hasattr(shape, "shapes"):
                if is_group:
                    # GROUPシェイプのa:offとa:chOffを取得
                    g_off_x, g_off_y, g_chOff_x, g_chOff_y = get_group_offsets(shape)
                    print(f"  デバッグ: GROUPシェイプを発見 (深度={depth}, 位置=({shape_left:.0f}, {shape_top:.0f}), 相対座標=({shape_relative_left:.0f}, {shape_relative_top:.0f}))")
                    if g_off_x is not None:
                        print(f"    GROUP座標系: off=({g_off_x:.0f}, {g_off_y:.0f}), chOff=({g_chOff_x:.0f}, {g_chOff_y:.0f})")
                    # GROUPシェイプ内の子シェイプを処理（a:offとa:chOffを渡す）
                    for child_shape in shape.shapes:
                        process_shape(child_shape, shape_left, shape_top, depth + 1, 
                                     group_off_x=g_off_x, group_off_y=g_off_y, 
                                     group_chOff_x=g_chOff_x, group_chOff_y=g_chOff_y,
                                     parent_group_shape=current_group_shape)
                else:
                    # 通常のシェイプの場合は、親のGROUP座標系情報をそのまま渡す
                    for child_shape in shape.shapes:
                        process_shape(child_shape, shape_left, shape_top, depth + 1,
                                     group_off_x=group_off_x, group_off_y=group_off_y,
                                     group_chOff_x=group_chOff_x, group_chOff_y=group_chOff_y,
                                     parent_group_shape=current_group_shape)
        except:
            pass
    
    # すべてのシェイプを処理
    for shape in slide.shapes:
        process_shape(shape, parent_left=0, parent_top=0, depth=0, 
                     group_off_x=None, group_off_y=None, group_chOff_x=None, group_chOff_y=None,
                     parent_group_shape=None)
    
    # デバッグ: スライドサイズを出力
    center_x = slide_width / 2
    center_y = slide_height / 2
    print(f"  デバッグ: スライドサイズ: 幅={slide_width:.0f}, 高さ={slide_height:.0f}, 中央=({center_x:.0f}, {center_y:.0f})")
    
    # デバッグ: すべての「点」の情報を出力（中央基準の-100～100座標系で表示）
    if all_ten_shapes:
        print(f"  デバッグ: スライド内で見つかった「点」のみの総数: {len(all_ten_shapes)}個")
        for i, ten_info in enumerate(all_ten_shapes):
            right = ten_info.get('right', ten_info['left'] + ten_info['width'])
            # 座標をスライド範囲内に制限（範囲外の場合は警告を表示）
            right_clamped = max(0, min(right, slide_width))
            top_clamped = max(0, min(ten_info['top'], slide_height))
            is_out_of_bounds = (right > slide_width or ten_info['top'] > slide_height or ten_info['left'] < 0 or ten_info['top'] < 0)
            
            # 中央基準の-100～100座標系に変換（制限後の座標を使用）
            # X座標: 右端が中央より右なら正、左なら負
            # スライドの右端が100、左端が-100になるように正規化
            x_normalized = ((right_clamped - center_x) / (slide_width / 2)) * 100 if slide_width > 0 else 0
            # Y座標: 上端が中央より下なら正、上なら負（PowerPointは上から下が正方向）
            # スライドの下端が100、上端が-100になるように正規化
            y_normalized = ((top_clamped - center_y) / (slide_height / 2)) * 100 if slide_height > 0 else 0
            
            right_ratio = (right / slide_width * 100) if slide_width > 0 else 0
            print(f"  デバッグ: 「点」のみ{i+1}: 右端位置={right:.0f} (右端比率={right_ratio:.1f}%)")
            if is_out_of_bounds:
                print(f"    ⚠️ 警告: 座標がスライド範囲外です（スライド幅={slide_width:.0f}, 高さ={slide_height:.0f}）")
                print(f"    制限後座標: 右端={right_clamped:.0f}, 上端={top_clamped:.0f}")
            print(f"    正規化座標: X={x_normalized:.1f}, Y={y_normalized:.1f} (中央(0,0)基準、-100～100)")
            print(f"    元座標: 左端={ten_info['left']:.0f}, 右端={right:.0f}, 上端={ten_info['top']:.0f}, テキスト='{ten_info.get('text', '')}'")
    
    # 「点」のみのテキストは問題数と同じだけ（各スライドに2個）のはず
    # フォントサイズ8ptの「点」のテキストボックスの位置を見つけ、それが相対的に左か右かで振り分ける
    # 中央（slide_width / 2）で分類する
    center_x = slide_width / 2
    left_tens = []
    right_tens = []
    
    for ten in all_ten_shapes:
        # テキストボックスの右端位置を取得（位置判定に使用）
        ten_right = ten.get('right', ten['left'] + ten['width'])
        ten_right_ratio = (ten_right / slide_width * 100) if slide_width > 0 else 0
        # テキストボックスの中心位置も計算（参考用）
        ten_center = ten['left'] + ten['width'] / 2
        ten_center_ratio = (ten_center / slide_width * 100) if slide_width > 0 else 0
        # 右端が中央より左にある場合は左側、中央より右にある場合は右側
        if ten_right < center_x:
            left_tens.append(ten)
            print(f"  デバッグ: 左側に分類（右端: {ten_right:.0f}, 右端比率: {ten_right_ratio:.1f}%, 中心: {ten_center:.0f}, 左端: {ten['left']:.0f}）")
        else:
            right_tens.append(ten)
            print(f"  デバッグ: 右側に分類（右端: {ten_right:.0f}, 右端比率: {ten_right_ratio:.1f}%, 中心: {ten_center:.0f}, 左端: {ten['left']:.0f}）")
    
    print(f"  デバッグ: 中央で分類（右端基準）: 左側={len(left_tens)}個, 右側={len(right_tens)}個, 中央={center_x:.0f}")
    
    # 指定された領域の「点」を選択
    if position == 'left':
        ten_shapes = [{'shape': ten['shape'], 'left': ten['left'], 'top': ten['top'], 
                       'width': ten['width'], 'height': ten['height'],
                       'parent_group_shape': ten.get('parent_group_shape')} for ten in left_tens]
    else:
        ten_shapes = [{'shape': ten['shape'], 'left': ten['left'], 'top': ten['top'], 
                       'width': ten['width'], 'height': ten['height'],
                       'parent_group_shape': ten.get('parent_group_shape')} for ten in right_tens]
    
    # 「点」のシェイプが見つからなかった場合
    if not ten_shapes:
        print(f"  デバッグ: 「点」が見つかりませんでした ({position}側)")
        return None
    
    # 各「点」のシェイプについて、その上に横向きの直線とテキストボックスがあるか確認
    for ten_info in ten_shapes:
        ten_shape = ten_info['shape']
        ten_top = ten_info['top']
        ten_left = ten_info['left']
        ten_width = ten_info['width']
        ten_height = ten_info['height']
        parent_group = ten_info.get('parent_group_shape')
        
        # まず、GROUPシェイプ内に「点」「横棒」「テキストボックス」がすべて含まれているか確認
        if parent_group is not None and hasattr(parent_group, 'shapes'):
            print(f"  デバッグ: 「点」がGROUPシェイプ内にあります。GROUP内の要素を確認します")
            
            # GROUPシェイプ内の要素を確認
            has_ten = False  # 「点」があるか
            has_line = False  # 横向きの直線があるか
            answer_textbox = None  # 解答欄のテキストボックス
            
            def check_group_contents(shape, group_off_x=None, group_chOff_x=None, group_off_y=None, group_chOff_y=None):
                """GROUPシェイプ内の要素を確認"""
                nonlocal has_ten, has_line, answer_textbox
                try:
                    is_group = hasattr(shape, 'shape_type') and shape.shape_type == 6
                    shape_relative_left = 0
                    shape_relative_top = 0
                    
                    if hasattr(shape, 'left') and hasattr(shape, 'top'):
                        shape_relative_left = shape.left
                        shape_relative_top = shape.top
                        
                        if group_off_x is not None and group_chOff_x is not None:
                            shape_left = shape_relative_left + group_off_x - group_chOff_x
                            shape_top = shape_relative_top + group_off_y - group_chOff_y
                        else:
                            shape_left = shape_relative_left
                            shape_top = shape_relative_top
                    else:
                        shape_left = 0
                        shape_top = 0
                    
                    # 「点」かどうか確認
                    text = extract_text_from_shape(shape)
                    if text.strip() == "点" and len(text.strip()) == 1:
                        font_size = get_font_size_from_shape(shape)
                        if font_size == 8.0:
                            has_ten = True
                            print(f"    デバッグ: GROUP内に「点」を発見 (位置: {shape_left:.0f}, {shape_top:.0f})")
                    
                    # 横向きの直線かどうか確認
                    if hasattr(shape, 'width') and hasattr(shape, 'height'):
                        if shape.width > shape.height * 2:  # 幅が高さの2倍以上
                            has_line = True
                            print(f"    デバッグ: GROUP内に横向きの直線を発見 (位置: {shape_left:.0f}, {shape_top:.0f}, サイズ: {shape.width:.0f}x{shape.height:.0f})")
                        elif hasattr(shape, 'shape_type') and shape.shape_type == 1:  # LINE
                            has_line = True
                            print(f"    デバッグ: GROUP内に横向きの直線を発見 (位置: {shape_left:.0f}, {shape_top:.0f}, タイプ: LINE)")
                    
                    # テキストボックスかどうか確認（「点」以外）
                    if hasattr(shape, 'text_frame') or (hasattr(shape, 'shape_type') and shape.shape_type == 17):
                        text = extract_text_from_shape(shape)
                        if text.strip() != "点":
                            # 解答欄のテキストボックス候補（最も下にあるものを選択）
                            if answer_textbox is None:
                                answer_textbox = shape
                                print(f"    デバッグ: GROUP内にテキストボックス候補を発見 (位置: {shape_left:.0f}, {shape_top:.0f}, テキスト: '{text[:20]}')")
                            else:
                                # 既存の候補と比較して、より下にあるものを選択
                                existing_top = answer_textbox.top if hasattr(answer_textbox, 'top') else 0
                                if shape_top > existing_top:
                                    answer_textbox = shape
                                    print(f"    デバッグ: GROUP内により下のテキストボックス候補を発見 (位置: {shape_left:.0f}, {shape_top:.0f}, テキスト: '{text[:20]}')")
                    
                    # GROUPシェイプ内の子シェイプも確認
                    if hasattr(shape, "shapes"):
                        if is_group:
                            g_off_x, g_off_y, g_chOff_x, g_chOff_y = get_group_offsets(shape)
                            for child_shape in shape.shapes:
                                check_group_contents(child_shape, g_off_x, g_chOff_x, g_off_y, g_chOff_y)
                        else:
                            for child_shape in shape.shapes:
                                check_group_contents(child_shape, group_off_x, group_chOff_x, group_off_y, group_chOff_y)
                except:
                    pass
            
            # GROUPシェイプのa:offとa:chOffを取得
            g_off_x, g_off_y, g_chOff_x, g_chOff_y = get_group_offsets(parent_group)
            for child_shape in parent_group.shapes:
                check_group_contents(child_shape, g_off_x, g_chOff_x, g_off_y, g_chOff_y)
            
            # 「点」「横棒」「テキストボックス」がすべて揃っていれば、そのテキストボックスが解答欄
            if has_ten and has_line and answer_textbox is not None:
                print(f"  デバッグ: GROUP内に「点」「横棒」「テキストボックス」がすべて揃っています。解答欄として確定します")
                return answer_textbox
            else:
                print(f"  デバッグ: GROUP内の要素確認: 点={has_ten}, 横棒={has_line}, テキストボックス={answer_textbox is not None}")
        
        # GROUPシェイプ内にすべて揃っていない場合は、従来の方法で探す
        # 「点」の上にあるシェイプを探す（上下方向のマージンを考慮）
        # 左側の場合は探索範囲を広めにする
        if position == 'left':
            search_top = ten_top - slide_height * 0.4  # 「点」の上40%の範囲を探索
        else:
            search_top = ten_top - slide_height * 0.3  # 「点」の上30%の範囲を探索
        search_bottom = ten_top
        
        # 「点」の上にあるシェイプをすべて収集
        candidate_shapes = []  # (shape, left, top, width, height, shape_type)
        
        def collect_above_shapes(shape, parent_left=0, parent_top=0, group_off_x=None, group_chOff_x=None, group_off_y=None, group_chOff_y=None):
            """「点」の上にあるシェイプを収集（GROUPシェイプ内の座標変換に対応）"""
            try:
                is_group = hasattr(shape, 'shape_type') and shape.shape_type == 6  # 6 = MSO_SHAPE_TYPE.GROUP
                shape_relative_left = 0
                shape_relative_top = 0
                
                if hasattr(shape, 'left') and hasattr(shape, 'top'):
                    shape_relative_left = shape.left
                    shape_relative_top = shape.top
                    
                    # GROUPシェイプ内の子シェイプの場合、a:offとa:chOffを考慮して絶対座標に変換
                    if group_off_x is not None and group_chOff_x is not None:
                        # absolute_x = child_shape.left + group_off_x - group_chOff_x
                        # absolute_y = child_shape.top + group_off_y - group_chOff_y
                        shape_left = shape_relative_left + group_off_x - group_chOff_x
                        shape_top = shape_relative_top + group_off_y - group_chOff_y
                    else:
                        # 通常のシェイプまたはGROUPシェイプ自体の場合
                        shape_left = parent_left + shape_relative_left
                        shape_top = parent_top + shape_relative_top
                else:
                    shape_left = parent_left
                    shape_top = parent_top
                
                # 左右の領域を確認
                # 左側: 0 ～ slide_width * 0.62 (中央より少し左まで)
                # 右側: slide_width * 0.38 ～ slide_width
                area = None
                if shape_left < left_area_end:
                    area = 'left'
                elif shape_left >= right_area_start:
                    area = 'right'
                
                # 指定された領域で、「点」の上にあるシェイプを収集
                if area == position and search_top <= shape_top <= search_bottom:
                    shape_width = shape.width if hasattr(shape, 'width') else 0
                    shape_height = shape.height if hasattr(shape, 'height') else 0
                    shape_type = shape.shape_type if hasattr(shape, 'shape_type') else None
                    
                    candidate_shapes.append({
                        'shape': shape,
                        'left': shape_left,
                        'top': shape_top,
                        'width': shape_width,
                        'height': shape_height,
                        'shape_type': shape_type
                    })
                
                # GROUPシェイプの場合は子シェイプも探索
                if hasattr(shape, "shapes"):
                    if is_group:
                        # GROUPシェイプのa:offとa:chOffを取得
                        g_off_x, g_off_y, g_chOff_x, g_chOff_y = get_group_offsets(shape)
                        # GROUPシェイプ内の子シェイプを処理（a:offとa:chOffを渡す）
                        for child_shape in shape.shapes:
                            collect_above_shapes(child_shape, shape_left, shape_top,
                                                group_off_x=g_off_x, group_off_y=g_off_y,
                                                group_chOff_x=g_chOff_x, group_chOff_y=g_chOff_y)
                    else:
                        # 通常のシェイプの場合は、親のGROUP座標系情報をそのまま渡す
                        for child_shape in shape.shapes:
                            collect_above_shapes(child_shape, shape_left, shape_top,
                                                group_off_x=group_off_x, group_off_y=group_off_y,
                                                group_chOff_x=group_chOff_x, group_chOff_y=group_chOff_y)
            except:
                pass
        
        # すべてのシェイプを再探索
        for shape in slide.shapes:
            collect_above_shapes(shape, parent_left=0, parent_top=0,
                               group_off_x=None, group_off_y=None, group_chOff_x=None, group_chOff_y=None)
        
        # 横向きの直線を探す（幅が高さより大きい、またはシェイプタイプが線）
        line_shapes = []
        for candidate in candidate_shapes:
            # 横向きの直線かどうかを判定
            is_horizontal_line = False
            if candidate['width'] > candidate['height'] * 2:  # 幅が高さの2倍以上
                is_horizontal_line = True
            elif candidate['shape_type'] == 1:  # MSO_SHAPE_TYPE.LINE
                is_horizontal_line = True
            
            if is_horizontal_line:
                line_shapes.append(candidate)
        
        print(f"  デバッグ: 「点」の上に横向きの直線 {len(line_shapes)}個を発見 ({position}側)")
        
        # 各直線について、その上にあるテキストボックスを探す
        for line_info in line_shapes:
            line_top = line_info['top']
            line_left = line_info['left']
            
            # 直線の上にあるテキストボックスを探す
            textbox_candidates = []
            for candidate in candidate_shapes:
                # テキストボックスかどうかを判定
                is_textbox = False
                if hasattr(candidate['shape'], 'text_frame'):
                    is_textbox = True
                elif candidate['shape_type'] == 17:  # MSO_SHAPE_TYPE.TEXT_BOX
                    is_textbox = True
                
                if is_textbox:
                    # 「点」のテキストボックスは除外する
                    text = extract_text_from_shape(candidate['shape'])
                    if text == "点" or text.strip() == "点":
                        continue
                    
                    # テキストボックスが直線の上にあるか確認（位置が近い）
                    if candidate['top'] < line_top:
                        textbox_right = candidate['left'] + candidate['width']
                        textbox_center = candidate['left'] + candidate['width'] / 2
                        
                        # 左右の位置も近いか確認（同じ領域内）
                        if position == 'left':
                            # 左側の場合は、テキストボックスの右端が中央付近にあることを確認
                            # または、テキストボックスの中心が左側領域内にあることを確認
                            # 左端が長いテキストボックスの場合、右端の条件を緩和
                            is_in_left_area = (candidate['left'] < slide_width / 2)  # 左端が中央より左
                            right_end_near_center = (slide_width * 0.4 <= textbox_right <= slide_width * 0.65)  # 右端が中央付近
                            center_in_left_area = (textbox_center < slide_width / 2)  # 中心が左側
                            
                            # 左側のテキストボックス: 左端が左側領域内、かつ右端が中央付近、または中心が左側
                            if is_in_left_area and (right_end_near_center or center_in_left_area):
                                # 直線との位置関係も確認（X座標が近い）
                                distance_threshold = slide_width * 0.4  # 左側は範囲を広めに
                                if abs(textbox_center - line_left) < distance_threshold or abs(candidate['left'] - line_left) < distance_threshold:
                                    textbox_candidates.append(candidate)
                                    print(f"  デバッグ: テキストボックス候補を発見 ({position}側, 左端: {candidate['left']:.0f}, 右端: {textbox_right:.0f}, 中心: {textbox_center:.0f}, 上端: {candidate['top']:.0f}, テキスト: '{text[:20]}')")
                            else:
                                print(f"  デバッグ: テキストボックスを除外 ({position}側, 左端: {candidate['left']:.0f}, 右端: {textbox_right:.0f}, 中心: {textbox_center:.0f}, 条件: 左端左側={is_in_left_area}, 右端中央付近={right_end_near_center}, 中心左側={center_in_left_area})")
                        else:
                            # 右側の場合は、従来通り距離で判定
                            distance_threshold = slide_width * 0.3
                            if abs(candidate['left'] - line_left) < distance_threshold:
                                textbox_candidates.append(candidate)
                                print(f"  デバッグ: テキストボックス候補を発見 ({position}側, 位置: {candidate['left']:.0f}, {candidate['top']:.0f}, テキスト: '{text[:20]}')")
            
            # 最も下にあるテキストボックスを選択（直線に最も近い）
            if textbox_candidates:
                textbox_info = max(textbox_candidates, key=lambda x: x['top'])
                print(f"  デバッグ: 解答欄のテキストボックスを発見 ({position}側)")
                return textbox_info['shape']
            else:
                print(f"  デバッグ: テキストボックス候補が見つかりませんでした ({position}側, 直線位置: {line_left:.0f}, {line_top:.0f})")
    
    return None


def add_answer_textbox(slide, answer_text, prs, position='right'):
    """
    スライドの下部に解答テキストボックスを追加、または既存の解答欄に入力（左右どちらかの領域）
    
    解答欄が存在する場合はそこに入力し、存在しない場合は新規にテキストボックスを追加する。
    
    Args:
        slide: 対象のスライドオブジェクト
        answer_text: 解答文字列（例: "38.608㎠"）
        prs: Presentationオブジェクト（スライドサイズ取得用）
        position: 'left' または 'right'（左右どちらの領域に配置するか）
    """
    # 既存の解答欄のテキストボックスを探す
    existing_textbox = find_answer_textbox(slide, prs, position)
    
    if existing_textbox is not None:
        # 既存の解答欄が見つかった場合、そこに解答を入力
        try:
            if hasattr(existing_textbox, 'text_frame'):
                text_frame = existing_textbox.text_frame
            else:
                # text_frameがない場合、テキストプロパティを直接設定
                existing_textbox.text = answer_text
                print(f"  デバッグ: 既存の解答欄に解答を入力 ({position}側): {answer_text}")
                return
            
            # 解答を設定（既存のテキストは完全に置き換える）
            text_frame.text = answer_text
            
            # 段落の設定（右寄せ）
            if len(text_frame.paragraphs) > 0:
                paragraph = text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.RIGHT
                
                # フォント設定（12pt、游明朝）
                font = paragraph.font
                font.size = Pt(12)  # 12pt
                font.bold = False  # 太字なし
                font.color.rgb = RGBColor(0, 0, 0)  # 黒色
                
                # 游明朝フォントを設定
                try:
                    font.name = "游明朝"
                except:
                    try:
                        font.name = "Yu Mincho"
                    except:
                        # 游明朝が利用できない場合はデフォルトフォントを使用
                        pass
            
            print(f"  デバッグ: 既存の解答欄に解答を入力 ({position}側): {answer_text}")
            return
        except Exception as e:
            print(f"  デバッグ: 既存の解答欄への入力に失敗: {e}")
            # エラーが発生した場合は、新規にテキストボックスを追加する処理に続く
    
    # 既存の解答欄が見つからなかった場合、新規にテキストボックスを追加
    # スライドのサイズを取得（プレゼンテーションから取得）
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # テキストボックスのサイズと位置を設定
    textbox_width = Inches(3)
    textbox_height = Inches(0.6)
    margin = Inches(0.5)  # 端からのマージン
    top = slide_height * 0.82  # 下部（82%の位置）
    
    # 左右どちらの領域に配置するかで位置を決定
    if position == 'left':
        # 左側の領域の右端に配置
        area_width = slide_width / 2
        left = area_width - textbox_width - margin
    else:  # position == 'right'
        # 右側の領域の右端に配置
        area_start = slide_width / 2
        left = slide_width - textbox_width - margin
    
    # テキストボックスを作成
    textbox = slide.shapes.add_textbox(left, top, textbox_width, textbox_height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = False
    
    # テキストを設定
    text_frame.text = answer_text
    
    # 段落の設定（右寄せ）
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.RIGHT
    
    # フォント設定（12pt、游明朝）
    font = paragraph.font
    font.size = Pt(12)  # 12pt
    font.bold = False  # 太字なし
    font.color.rgb = RGBColor(0, 0, 0)  # 黒色
    
    # 游明朝フォントを設定
    try:
        font.name = "游明朝"
    except:
        try:
            font.name = "Yu Mincho"
        except:
            # 游明朝が利用できない場合はデフォルトフォントを使用
            pass


def convert_pptx(input_path, output_path=None):
    """
    PPTXファイルを変換する
    
    Args:
        input_path: 入力PPTXファイルのパス
        output_path: 出力PPTXファイルのパス（Noneの場合は自動生成）
    """
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"入力ファイルが見つかりません: {input_path}")
    
    if not input_path.lower().endswith('.pptx'):
        raise ValueError("PPTXファイルを指定してください")
    
    # 出力パスが指定されていない場合は自動生成
    if output_path is None:
        base, ext = os.path.splitext(input_path)
        output_path = f"{base}_converted{ext}"
    
    # PPTXを読み込む
    prs = Presentation(input_path)
    
    if len(prs.slides) < 2:
        raise ValueError("スライドが2枚未満です。表紙と解答ページが必要です。")
    
    # スライド構成を判定
    # 0枚目 → 表紙
    # 最後 → 解答ページ
    # それ以外 → 問題ページ群
    # 注意: prs.slidesはスライスをサポートしていないため、リストに変換する
    slides_list = list(prs.slides)
    cover_slide = slides_list[0]
    answer_slide = slides_list[-1]
    question_slides = slides_list[1:-1]
    
    if len(question_slides) == 0:
        raise ValueError("問題ページが見つかりません。")
    
    print(f"スライド構成:")
    print(f"  表紙: 1枚")
    print(f"  問題ページ: {len(question_slides)}枚")
    print(f"  解答ページ: 1枚")
    
    # 解答ページから解答を抽出
    answers = extract_answers(answer_slide)
    print(f"\n抽出された解答: {len(answers)}件")
    for q_num, answer in sorted(answers.items()):
        print(f"  大問 {q_num}: {answer}")
    
    # 元のファイルをコピーしてから操作（画像やコメントも保持）
    # 出力ファイルとしてコピーを作成
    shutil.copy2(input_path, output_path)
    
    # コピーしたファイルを読み込む
    new_prs = Presentation(output_path)
    
    # スライドリストを取得
    new_slides_list = list(new_prs.slides)
    
    # 各問題ページから大問番号を抽出して照合
    answer_numbers = set(answers.keys())
    slide_question_numbers = match_question_numbers(question_slides, new_prs, answer_numbers)
    
    # 問題ページを複製して解答入りページを作成
    # 問題ページの後に挿入するため、前から処理
    new_question_slides_with_answers = []
    
    for idx in range(len(question_slides)):  # 0から始まるインデックス
        # 元のスライドのインデックスを取得（new_prs内での位置）
        # 0: 表紙, 1～len(question_slides): 問題ページ, 最後: 解答ページ
        source_slide_index = idx + 1  # 問題ページのインデックス（1から始まる）
        
        # 元のスライドを取得
        source_slide = new_slides_list[source_slide_index]
        
        # 照合された大問番号を取得（左右両方）
        left_num, right_num = slide_question_numbers[idx]
        
        # スライドを完全複製
        new_slide = duplicate_slide_complete(new_prs, source_slide)
        
        # 左側の解答を追加
        if left_num is not None:
            answer_text = answers.get(left_num, "（未設定）")
            if answer_text == "（未設定）":
                print(f"警告: 大問{left_num}の解答が見つかりませんでした。")
            else:
                add_answer_textbox(new_slide, answer_text, new_prs, position='left')
        
        # 右側の解答を追加
        if right_num is not None:
            answer_text = answers.get(right_num, "（未設定）")
            if answer_text == "（未設定）":
                print(f"警告: 大問{right_num}の解答が見つかりませんでした。")
            else:
                add_answer_textbox(new_slide, answer_text, new_prs, position='right')
        
        new_question_slides_with_answers.append(new_slide)
    
    # スライドの順序を調整: 表紙、問題ページ群、解答入り問題ページ群、解答ページ
    # python-pptxでは直接順序を変更できないため、XMLレベルで順序を調整
    # スライドIDリスト（_sldIdLst）の順序を変更
    sldIdLst = new_prs.slides._sldIdLst
    
    # 現在のスライドID要素を取得
    current_slide_ids = list(sldIdLst)
    
    # 順序: 表紙(0), 問題ページ(1～len), 解答入り問題ページ(追加分), 解答ページ(最後)
    # 新しい順序でスライドID要素を再配置
    new_order = []
    
    # 1. 表紙（最初のスライド）
    new_order.append(current_slide_ids[0])
    
    # 2. 問題ページ群（1～len(question_slides)）
    for i in range(1, len(question_slides) + 1):
        new_order.append(current_slide_ids[i])
    
    # 3. 解答入り問題ページ群（最後に追加されたスライド）
    # 追加されたスライドは最後に配置されるため、問題ページの後から解答ページの前まで
    for i in range(len(question_slides) + 1, len(current_slide_ids) - 1):
        new_order.append(current_slide_ids[i])
    
    # 4. 解答ページ（最後のスライド）
    new_order.append(current_slide_ids[-1])
    
    # スライドIDリストをクリアして再構築
    sldIdLst.clear()
    for slide_id_elem in new_order:
        sldIdLst.append(slide_id_elem)
    
    # 新しいPPTXを保存
    new_prs.save(output_path)
    print(f"\n変換完了: {output_path}")
    print(f"  総スライド数: {len(new_prs.slides)}枚")
    print(f"  (表紙: 1枚, 問題ページ: {len(question_slides)}枚, 解答入り問題ページ: {len(new_question_slides_with_answers)}枚, 解答ページ: 1枚)")


def main():
    """コマンドライン実行用のメイン関数"""
    if len(sys.argv) < 2:
        print("使用方法: python convert_pptx.py <入力PPTXファイル> [出力PPTXファイル]")
        print("例: python convert_pptx.py input.pptx output.pptx")
        sys.exit(1)
    
    input_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None
    
    try:
        convert_pptx(input_path, output_path)
    except Exception as e:
        print(f"エラー: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()





